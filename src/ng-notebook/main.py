import os
from typing import List
from fastapi import FastAPI, UploadFile, File
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import uvicorn
from langchain_community.llms import Ollama
from langchain_community.embeddings import OllamaEmbeddings
from langchain_community.vectorstores import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import ConversationBufferMemory
import streamlit as st
import tempfile
import pandas as pd
from pptx import Presentation
import PyPDF2
import magic
from langchain.prompts import PromptTemplate
import shutil
import chromadb
from chromadb.config import Settings
from langchain_community.vectorstores.utils import filter_complex_metadata

# Initialize FastAPI app
app = FastAPI()

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Track uploaded documents
uploaded_documents = set()

# Initialize Ollama
llm = Ollama(model="llama3.3:latest")
embeddings = OllamaEmbeddings(model="all-minilm")

# Ensure chroma_db directory exists and has proper permissions
CHROMA_DB_DIR = "./chroma_db"
if os.path.exists(CHROMA_DB_DIR):
    # Remove existing directory if it exists
    shutil.rmtree(CHROMA_DB_DIR)
os.makedirs(CHROMA_DB_DIR, mode=0o777)

# Initialize Chroma client with proper settings
chroma_client = chromadb.PersistentClient(
    path=CHROMA_DB_DIR,
    settings=Settings(
        anonymized_telemetry=False,
        allow_reset=True,
        is_persistent=True
    )
)

# Create the collection
collection = chroma_client.create_collection(
    name="documents",
    metadata={"hnsw:space": "cosine"}
)

# Initialize Chroma as the vector store
vector_store = Chroma(
    client=chroma_client,
    embedding_function=embeddings,
    collection_name="documents",
    persist_directory=CHROMA_DB_DIR
)

# Initialize conversation memory
memory = ConversationBufferMemory(
    memory_key="chat_history",
    return_messages=True,
    output_key="answer"
)

# Create a more detailed prompt template for better document handling
template = """You are an AI assistant specialized in analyzing and explaining data from various document types. 
When dealing with different document types, pay special attention to:

For Excel/CSV files:
1. Numerical data and calculations
2. Column names and their relationships
3. Row indices and their context
4. Data types and formats
5. Sheet names and their relationships (for Excel)
6. Statistical information (mean, median, correlations)
7. Sample values and unique counts

For PDF documents:
1. Page numbers and their context
2. Document structure and sections
3. Headers and subheaders
4. Lists and bullet points
5. Tables and their content
6. Important figures and statistics
7. Key concepts and their relationships

For PowerPoint presentations:
1. Slide numbers and their sequence
2. Slide titles and subtitles
3. Bullet points and their hierarchy
4. Images and their captions
5. Speaker notes if available
6. Presentation flow and structure
7. Key messages and takeaways

Use the following pieces of context to answer the question at the end.
If you don't know the answer, just say that you don't know, don't try to make up an answer.

Context: {context}

Chat History: {chat_history}

Question: {question}

When answering:
1. For Excel/CSV data:
   - Reference specific columns, rows, or cells when relevant
   - Include numerical values in your response
   - Explain patterns or relationships in the data
   - Show calculations when asked
   - Mention sheet names for Excel files

2. For PDF documents:
   - Reference specific pages or sections
   - Maintain document structure in your response
   - Include relevant quotes or statistics
   - Explain relationships between concepts
   - Reference tables or figures when relevant

3. For PowerPoint presentations:
   - Reference specific slides
   - Maintain presentation flow
   - Include key points and their context
   - Explain relationships between slides
   - Reference visual elements when relevant

4. General guidelines:
   - Be specific and precise in your references
   - Provide context for your answers
   - Explain relationships and patterns
   - Include relevant statistics or numbers
   - Maintain document structure in your response

Answer:"""

# Create the prompt
prompt = PromptTemplate(
    input_variables=["context", "chat_history", "question"],
    template=template
)

# Initialize the chain with the custom prompt
chain = ConversationalRetrievalChain.from_llm(
    llm=llm,
    retriever=vector_store.as_retriever(
        search_type="similarity"
    ),
    memory=memory,
    return_source_documents=True,
    combine_docs_chain_kwargs={"prompt": prompt}
)

class ChatMessage(BaseModel):
    role: str
    content: str

class ChatRequest(BaseModel):
    message: str
    chat_history: List[ChatMessage]

def process_document(file_path: str, file_type: str) -> List[dict]:
    """Process different types of documents and return chunks with metadata."""
    chunks = []
    
    try:
        if file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            # Read all sheets from the Excel file
            excel_file = pd.ExcelFile(file_path)
            
            # Process each sheet
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                
                # Detect sub-tables by looking for empty rows
                empty_rows = df.index[df.isna().all(axis=1)].tolist()
                sub_tables = []
                start_idx = 0
                
                for empty_row in empty_rows:
                    if empty_row > start_idx:
                        sub_tables.append(df.iloc[start_idx:empty_row])
                    start_idx = empty_row + 1
                
                # Add the last sub-table if there is one
                if start_idx < len(df):
                    sub_tables.append(df.iloc[start_idx:])
                
                # If no sub-tables were found, treat the entire sheet as one table
                if not sub_tables:
                    sub_tables = [df]
                
                # Process each sub-table
                for table_idx, sub_table in enumerate(sub_tables):
                    # Clean the sub-table
                    sub_table = sub_table.dropna(how='all').dropna(axis=1, how='all')
                    
                    # Skip empty tables
                    if sub_table.empty:
                        continue
                    
                    # Add table-level information
                    table_info = f"Excel Sub-table {table_idx + 1} in sheet '{sheet_name}' Information:\n"
                    table_info += f"Total Rows: {len(sub_table)}\n"
                    table_info += f"Total Columns: {len(sub_table.columns)}\n"
                    table_info += f"Column Names: {', '.join(sub_table.columns)}\n"
                    table_info += f"Memory Usage: {sub_table.memory_usage(deep=True).sum() / 1024:.2f} KB\n"
                    
                    # Add data types information
                    table_info += "\nColumn Data Types:\n"
                    for col in sub_table.columns:
                        dtype = sub_table[col].dtype
                        if pd.api.types.is_datetime64_any_dtype(dtype):
                            table_info += f"{col}: datetime\n"
                        elif pd.api.types.is_numeric_dtype(dtype):
                            table_info += f"{col}: numeric ({dtype})\n"
                        else:
                            table_info += f"{col}: {dtype}\n"
                    
                    chunks.append({
                        "content": table_info,
                        "metadata": {
                            "source": file_path,
                            "type": "excel_subtable_info",
                            "sheet_name": sheet_name,
                            "subtable_index": str(table_idx + 1),
                            "row_count": str(len(sub_table)),
                            "column_count": str(len(sub_table.columns)),
                            "columns": str(list(sub_table.columns))
                        }
                    })
                    
                    # Process each column in the sub-table
                    for column in sub_table.columns:
                        # Get column statistics
                        col_stats = {}
                        
                        # Handle different data types
                        if pd.api.types.is_datetime64_any_dtype(sub_table[column]):
                            col_stats.update({
                                "min_date": sub_table[column].min(),
                                "max_date": sub_table[column].max(),
                                "unique_dates": sub_table[column].nunique(),
                                "null_count": sub_table[column].isnull().sum()
                            })
                        elif pd.api.types.is_numeric_dtype(sub_table[column]):
                            col_stats.update({
                                "mean": sub_table[column].mean(),
                                "median": sub_table[column].median(),
                                "std": sub_table[column].std(),
                                "min": sub_table[column].min(),
                                "max": sub_table[column].max(),
                                "unique_values": sub_table[column].nunique(),
                                "null_count": sub_table[column].isnull().sum()
                            })
                        else:
                            col_stats.update({
                                "unique_values": sub_table[column].nunique(),
                                "null_count": sub_table[column].isnull().sum(),
                                "most_common": sub_table[column].value_counts().head(5).to_dict()
                            })
                        
                        # Create a detailed description of the column
                        col_description = f"Column '{column}' in sub-table {table_idx + 1} of sheet '{sheet_name}':\n"
                        col_description += f"Data Type: {sub_table[column].dtype}\n"
                        
                        # Add type-specific statistics
                        if pd.api.types.is_datetime64_any_dtype(sub_table[column]):
                            col_description += f"Date Range: {col_stats['min_date']} to {col_stats['max_date']}\n"
                            col_description += f"Unique Dates: {col_stats['unique_dates']}\n"
                        elif pd.api.types.is_numeric_dtype(sub_table[column]):
                            col_description += f"Mean: {col_stats['mean']:.2f}\n"
                            col_description += f"Median: {col_stats['median']:.2f}\n"
                            col_description += f"Standard Deviation: {col_stats['std']:.2f}\n"
                            col_description += f"Range: {col_stats['min']:.2f} to {col_stats['max']:.2f}\n"
                        else:
                            col_description += "Most Common Values:\n"
                            for value, count in col_stats['most_common'].items():
                                col_description += f"- {value}: {count} occurrences\n"
                        
                        col_description += f"Null Values: {col_stats['null_count']}\n"
                        
                        # Add sample values (formatted appropriately)
                        sample_values = sub_table[column].dropna().head(5)
                        if pd.api.types.is_datetime64_any_dtype(sub_table[column]):
                            sample_values = sample_values.dt.strftime('%Y-%m-%d %H:%M:%S')
                        col_description += f"Sample Values: {sample_values.tolist()}\n"
                        
                        # Add column content (formatted appropriately)
                        col_content = sub_table[column].astype(str)
                        if pd.api.types.is_datetime64_any_dtype(sub_table[column]):
                            col_content = sub_table[column].dt.strftime('%Y-%m-%d %H:%M:%S')
                        col_description += f"\nColumn Content:\n{col_content.str.cat(sep='\n')}"
                        
                        chunks.append({
                            "content": col_description,
                            "metadata": {
                                "source": file_path,
                                "type": "excel_column",
                                "sheet_name": sheet_name,
                                "subtable_index": str(table_idx + 1),
                                "column_name": column,
                                "row_count": str(len(sub_table)),
                                "column_index": str(sub_table.columns.get_loc(column)),
                                "statistics": str(col_stats)
                            }
                        })
                    
                    # Add correlation information for numerical columns in this sub-table
                    numeric_cols = sub_table.select_dtypes(include=['int64', 'float64']).columns
                    if len(numeric_cols) > 1:
                        corr_matrix = sub_table[numeric_cols].corr()
                        corr_info = f"Correlation Matrix for Numerical Columns in sub-table {table_idx + 1} of sheet '{sheet_name}':\n"
                        corr_info += str(corr_matrix)
                        
                        chunks.append({
                            "content": corr_info,
                            "metadata": {
                                "source": file_path,
                                "type": "excel_correlations",
                                "sheet_name": sheet_name,
                                "subtable_index": str(table_idx + 1),
                                "numeric_columns": str(list(numeric_cols))
                            }
                        })
                    
                    # Add row-level information for the first 100 rows (to avoid too many chunks)
                    for idx, row in sub_table.head(100).iterrows():
                        row_data = {}
                        for col, val in row.items():
                            if pd.api.types.is_datetime64_any_dtype(sub_table[col]):
                                row_data[col] = val.strftime('%Y-%m-%d %H:%M:%S') if pd.notnull(val) else None
                            else:
                                row_data[col] = str(val) if pd.notnull(val) else None
                        
                        row_info = f"Row {idx + 1} in sub-table {table_idx + 1} of sheet '{sheet_name}':\n"
                        row_info += "\n".join(f"{col}: {val}" for col, val in row_data.items())
                        
                        chunks.append({
                            "content": row_info,
                            "metadata": {
                                "source": file_path,
                                "type": "excel_row",
                                "sheet_name": sheet_name,
                                "subtable_index": str(table_idx + 1),
                                "row_index": str(idx),
                                "columns": str(list(sub_table.columns))
                            }
                        })
        
        elif file_type == "application/pdf":
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                
                # Split text into chunks
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=200,
                    length_function=len,
                    separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
                )
                
                # Create chunks with metadata
                for i, chunk in enumerate(text_splitter.split_text(text)):
                    chunks.append({
                        "content": chunk,
                        "metadata": {
                            "source": file_path,
                            "type": "pdf",
                            "page": str(i + 1),
                            "total_pages": str(len(pdf_reader.pages))
                        }
                    })
        
        elif file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            # Split text into chunks
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
                separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
            )
            
            # Create chunks with metadata
            for i, chunk in enumerate(text_splitter.split_text(text)):
                chunks.append({
                    "content": chunk,
                    "metadata": {
                        "source": file_path,
                        "type": "ppt",
                        "slide": str(i + 1),
                        "total_slides": str(len(prs.slides))
                    }
                })
        
        elif file_type in ["text/csv", "application/vnd.ms-excel"]:
            df = pd.read_csv(file_path)
            text = df.to_string()
            
            # Split text into chunks
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
                separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
            )
            
            # Create chunks with metadata
            for i, chunk in enumerate(text_splitter.split_text(text)):
                chunks.append({
                    "content": chunk,
                    "metadata": {
                        "source": file_path,
                        "type": "csv",
                        "row_count": str(len(df)),
                        "column_count": str(len(df.columns))
                    }
                })
        
        else:
            # Handle unknown file types by reading as text
            with open(file_path, 'r', encoding='utf-8') as file:
                text = file.read()
            
            # Split text into chunks
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                length_function=len,
                separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
            )
            
            # Create chunks with metadata
            for i, chunk in enumerate(text_splitter.split_text(text)):
                chunks.append({
                    "content": chunk,
                    "metadata": {
                        "source": file_path,
                        "type": "text",
                        "chunk_index": str(i + 1)
                    }
                })
    
    except Exception as e:
        # If any error occurs, create a single chunk with error information
        chunks.append({
            "content": f"Error processing file: {str(e)}",
            "metadata": {
                "source": file_path,
                "type": "error",
                "error_message": str(e)
            }
        })
    
    return chunks

def clean_metadata(metadata: dict) -> dict:
    """Clean metadata to ensure all values are simple types."""
    cleaned = {}
    for key, value in metadata.items():
        if isinstance(value, (str, int, float, bool)) or value is None:
            cleaned[key] = value
        elif isinstance(value, list):
            cleaned[key] = ",".join(str(x) for x in value)
        else:
            cleaned[key] = str(value)
    return cleaned

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """Handle file upload and process it for RAG."""
    global vector_store, uploaded_documents, collection
    
    # Check if file has already been uploaded
    if file.filename in uploaded_documents:
        return {"message": "File already uploaded"}
    
    # Create a temporary file
    with tempfile.NamedTemporaryFile(delete=False) as temp_file:
        content = await file.read()
        temp_file.write(content)
        temp_file_path = temp_file.name
    
    try:
        # Detect file type
        file_type = magic.from_file(temp_file_path, mime=True)
        
        # Process the document
        chunks = process_document(temp_file_path, file_type)
        
        # Add file metadata to each chunk and clean metadata
        for chunk in chunks:
            chunk["metadata"].update({
                "source": file.filename,
                "file_type": file_type
            })
            # Clean metadata to ensure all values are simple types
            chunk["metadata"] = clean_metadata(chunk["metadata"])
        
        # Add texts to vector store
        vector_store.add_texts(
            [chunk["content"] for chunk in chunks],
            metadatas=[chunk["metadata"] for chunk in chunks]
        )
        
        # Add to uploaded documents set
        uploaded_documents.add(file.filename)
        
        return {"message": "File processed successfully"}
    
    finally:
        # Clean up temporary file
        os.unlink(temp_file_path)

@app.post("/chat")
async def chat(request: ChatRequest):
    """Handle chat requests with RAG."""
    if vector_store is None:
        return {"response": "Please upload a document first."}
    
    # Get response
    response = chain({"question": request.message})
    
    return {
        "response": response["answer"],
        "sources": [doc.page_content for doc in response["source_documents"]]
    }

@app.post("/clear-db")
async def clear_database():
    """Clear the Chroma database."""
    global vector_store, uploaded_documents, chroma_client, collection
    try:
        # First, try to delete the collection
        try:
            chroma_client.delete_collection("documents")
        except Exception as e:
            print(f"Error deleting collection: {str(e)}")
        
        # Then delete the Chroma database directory
        if os.path.exists(CHROMA_DB_DIR):
            try:
                shutil.rmtree(CHROMA_DB_DIR)
            except Exception as e:
                print(f"Error removing directory: {str(e)}")
                # Try to remove files individually
                for root, dirs, files in os.walk(CHROMA_DB_DIR, topdown=False):
                    for name in files:
                        try:
                            os.remove(os.path.join(root, name))
                        except Exception as e:
                            print(f"Error removing file {name}: {str(e)}")
                    for name in dirs:
                        try:
                            os.rmdir(os.path.join(root, name))
                        except Exception as e:
                            print(f"Error removing directory {name}: {str(e)}")
                try:
                    os.rmdir(CHROMA_DB_DIR)
                except Exception as e:
                    print(f"Error removing root directory: {str(e)}")
        
        # Create fresh directory with proper permissions
        os.makedirs(CHROMA_DB_DIR, mode=0o777, exist_ok=True)
        
        # Reinitialize Chroma client
        chroma_client = chromadb.PersistentClient(
            path=CHROMA_DB_DIR,
            settings=Settings(
                anonymized_telemetry=False,
                allow_reset=True,
                is_persistent=True
            )
        )
        
        # Create a new collection
        collection = chroma_client.create_collection(
            name="documents",
            metadata={"hnsw:space": "cosine"}
        )
        
        # Reset the vector store
        vector_store = Chroma(
            client=chroma_client,
            embedding_function=embeddings,
            collection_name="documents",
            persist_directory=CHROMA_DB_DIR
        )
        
        # Clear uploaded documents set
        uploaded_documents.clear()
        
        return {"message": "Database cleared successfully"}
    except Exception as e:
        return {"error": f"Error clearing database: {str(e)}"}

@app.get("/collections")
async def get_collections():
    """Get information about Chroma collections."""
    try:
        if vector_store is None:
            return {"collections": [], "message": "No collections found"}
        
        # Get collection information
        collection = vector_store._collection
        count = collection.count()
        
        # Get unique metadata values
        metadata = collection.get()["metadatas"]
        unique_sources = set()
        unique_types = set()
        
        for meta in metadata:
            if meta:
                if "source" in meta:
                    unique_sources.add(meta["source"])
                if "type" in meta:
                    unique_types.add(meta["type"])
        
        return {
            "total_documents": count,
            "unique_sources": list(unique_sources),
            "document_types": list(unique_types),
            "collection_name": collection.name,
            "uploaded_files": list(uploaded_documents)
        }
    except Exception as e:
        return {"error": f"Error getting collections: {str(e)}"}

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000) 