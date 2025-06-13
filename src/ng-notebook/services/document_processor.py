import os
import pandas as pd
from pptx import Presentation
import PyPDF2
import magic
from typing import List, Dict
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import UnstructuredPowerPointLoader
from ..core.config import CHUNK_SIZE, CHUNK_OVERLAP

class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            length_function=len,
            separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
        )

    def process_document(self, file_path: str, file_type: str) -> List[Dict]:
        """Process different types of documents and return chunks with metadata."""
        chunks = []
        
        try:
            if file_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                chunks.extend(self._process_pptx(file_path))
            elif file_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
                chunks.extend(self._process_excel(file_path))
            elif file_type == "application/pdf":
                chunks.extend(self._process_pdf(file_path))
            else:
                raise ValueError(f"Unsupported file type: {file_type}")
                
            return chunks
        except Exception as e:
            raise Exception(f"Error processing document: {str(e)}")

    def _process_pptx(self, file_path: str) -> List[Dict]:
        chunks = []
        loader = UnstructuredPowerPointLoader(file_path)
        documents = loader.load()
        
        for doc in documents:
            slide_number = doc.metadata.get("page_number", "unknown")
            
            slide_info = f"PowerPoint Slide {slide_number} Information:\n"
            slide_info += f"Content Type: {doc.metadata.get('content_type', 'unknown')}\n"
            slide_info += f"Content:\n{doc.page_content}\n"
            
            chunks.append({
                "content": slide_info,
                "metadata": {
                    "source": file_path,
                    "type": "ppt_slide_info",
                    "slide_number": str(slide_number),
                    "content_type": doc.metadata.get("content_type", "unknown")
                }
            })
            
            content_chunks = self.text_splitter.split_text(doc.page_content)
            for i, chunk in enumerate(content_chunks):
                chunks.append({
                    "content": chunk,
                    "metadata": {
                        "source": file_path,
                        "type": "ppt_content",
                        "slide_number": str(slide_number),
                        "chunk_index": str(i + 1),
                        "total_chunks": str(len(content_chunks))
                    }
                })
        
        prs = Presentation(file_path)
        presentation_info = self._get_presentation_info(prs, file_path)
        chunks.append(presentation_info)
        
        return chunks

    def _process_excel(self, file_path: str) -> List[Dict]:
        chunks = []
        df = pd.read_excel(file_path)
        
        # Process each sheet
        for sheet_name in df.keys():
            sheet_df = df[sheet_name]
            
            # Add sheet information
            sheet_info = f"Excel Sheet Information:\n"
            sheet_info += f"Sheet Name: {sheet_name}\n"
            sheet_info += f"Rows: {len(sheet_df)}\n"
            sheet_info += f"Columns: {', '.join(sheet_df.columns)}\n"
            
            chunks.append({
                "content": sheet_info,
                "metadata": {
                    "source": file_path,
                    "type": "excel_sheet_info",
                    "sheet_name": sheet_name
                }
            })
            
            # Process data in chunks
            for i in range(0, len(sheet_df), 100):
                chunk_df = sheet_df.iloc[i:i+100]
                chunk_text = chunk_df.to_string()
                
                chunks.append({
                    "content": chunk_text,
                    "metadata": {
                        "source": file_path,
                        "type": "excel_data",
                        "sheet_name": sheet_name,
                        "start_row": str(i + 1),
                        "end_row": str(min(i + 100, len(sheet_df)))
                    }
                })
        
        return chunks

    def _process_pdf(self, file_path: str) -> List[Dict]:
        chunks = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                
                # Split text into chunks
                text_chunks = self.text_splitter.split_text(text)
                
                for i, chunk in enumerate(text_chunks):
                    chunks.append({
                        "content": chunk,
                        "metadata": {
                            "source": file_path,
                            "type": "pdf",
                            "page": str(page_num + 1),
                            "chunk": str(i + 1)
                        }
                    })
        
        return chunks

    def _get_presentation_info(self, prs: Presentation, file_path: str) -> Dict:
        presentation_info = f"PowerPoint Presentation Information:\n"
        presentation_info += f"Total Slides: {len(prs.slides)}\n"
        presentation_info += f"Slide Layouts Used: {', '.join(set(slide.slide_layout.name for slide in prs.slides))}\n"
        
        presentation_info += "\nSlide Sequence:\n"
        for i, slide in enumerate(prs.slides, 1):
            title = "No Title"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
                    break
            presentation_info += f"Slide {i}: {title}\n"
        
        return {
            "content": presentation_info,
            "metadata": {
                "source": file_path,
                "type": "ppt_presentation_info",
                "total_slides": str(len(prs.slides))
            }
        } 