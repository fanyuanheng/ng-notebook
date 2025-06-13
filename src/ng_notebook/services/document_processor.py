import os
import pandas as pd
from pptx import Presentation
import PyPDF2
import magic
from typing import List, Dict
import logging
import re
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import UnstructuredPowerPointLoader
from ..core.config import CHUNK_SIZE, CHUNK_OVERLAP
from .sqlite_store import SQLiteStore, sanitize_table_name
import sqlite3
from langchain.schema import Document

# Get logger
logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self, sqlite_store: SQLiteStore = None):
        logger.debug("Initializing DocumentProcessor")
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            length_function=len,
            separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
        )
        self.sqlite_store = sqlite_store or SQLiteStore()
        logger.debug("DocumentProcessor initialized successfully")

    def process_document(self, file_path: str, file_type: str) -> List[Document]:
        """Process a document into chunks."""
        logger.info(f"Processing document: {file_path} of type {file_type}")
        
        try:
            # Handle Excel files
            if file_type in [
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                "application/vnd.ms-excel",
                "application/x-excel",
                "application/x-msexcel",
                "application/vnd.ms-excel.sheet.macroEnabled.12",
                "application/vnd.ms-excel.sheet.binary.macroEnabled.12"
            ]:
                logger.info("Processing Excel file")
                # Store in SQLite first
                result = self.sqlite_store.add_excel_file(file_path)
                
                # Create summary chunks for vector store
                chunks = []
                with sqlite3.connect(self.sqlite_store.db_path) as conn:
                    cursor = conn.cursor()
                    cursor.execute("SELECT id FROM uploaded_files WHERE filename = ?", (os.path.basename(file_path),))
                    file_id = cursor.fetchone()[0]
                    
                    cursor.execute("SELECT sheet_name, row_count, column_count FROM excel_sheets WHERE file_id = ?", (file_id,))
                    sheets = cursor.fetchall()
                    
                    for sheet_name, row_count, column_count in sheets:
                        # Get table name
                        safe_sheet_name = sanitize_table_name(sheet_name)
                        table_name = f"excel_{file_id}_{safe_sheet_name}"
                        logger.debug(f"Using table name: {table_name}")
                        
                        try:
                            df = pd.read_sql(f"SELECT * FROM {table_name}", conn)
                            
                            # Create a summary of the sheet
                            summary = f"Sheet: {sheet_name}\n"
                            summary += f"Rows: {row_count}, Columns: {column_count}\n"
                            summary += f"Columns: {', '.join(df.columns)}\n"
                            summary += f"First few rows:\n{df.head().to_string()}"
                            
                            doc = Document(
                                page_content=summary,
                                metadata={
                                    "source": file_path,
                                    "file_type": file_type,
                                    "sheet_name": sheet_name,
                                    "row_count": row_count,
                                    "column_count": column_count
                                }
                            )
                            chunks.append(doc)
                        except Exception as e:
                            logger.error(f"Error processing sheet {sheet_name}: {str(e)}", exc_info=True)
                            # Continue with other sheets even if one fails
                            continue
                
                if not chunks:
                    raise Exception("No sheets were successfully processed")
                
                return chunks
                
            # Handle CSV files
            elif file_type in ["text/csv", "application/csv", "text/x-csv", "application/x-csv", "text/comma-separated-values"]:
                logger.info("Processing CSV file")
                # Store in SQLite first
                result = self.sqlite_store.add_csv_file(file_path)
                
                # Create summary chunk for vector store
                with sqlite3.connect(self.sqlite_store.db_path) as conn:
                    cursor = conn.cursor()
                    cursor.execute("SELECT id FROM uploaded_files WHERE filename = ?", (os.path.basename(file_path),))
                    file_id = cursor.fetchone()[0]
                    
                    # Get table name
                    safe_filename = sanitize_table_name(os.path.basename(file_path))
                    table_name = f"csv_{file_id}_{safe_filename}"
                    logger.debug(f"Using table name: {table_name}")
                    
                    try:
                        df = pd.read_sql(f"SELECT * FROM {table_name}", conn)
                        
                        # Create a summary of the CSV
                        summary = f"CSV File: {os.path.basename(file_path)}\n"
                        summary += f"Rows: {len(df)}, Columns: {len(df.columns)}\n"
                        summary += f"Columns: {', '.join(df.columns)}\n"
                        summary += f"First few rows:\n{df.head().to_string()}"
                        
                        doc = Document(
                            page_content=summary,
                            metadata={
                                "source": file_path,
                                "file_type": file_type,
                                "row_count": len(df),
                                "column_count": len(df.columns)
                            }
                        )
                        return [doc]
                    except Exception as e:
                        logger.error(f"Error processing CSV file: {str(e)}", exc_info=True)
                        raise
            
            # Handle PowerPoint files
            elif file_type in [
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "application/vnd.ms-powerpoint",
                "application/x-mspowerpoint"
            ]:
                logger.info("Processing PowerPoint file")
                chunks = self._process_pptx(file_path)
                return [Document(page_content=chunk["content"], metadata=chunk["metadata"]) for chunk in chunks]
            
            # Handle PDF files
            elif file_type in ["application/pdf", "application/x-pdf"]:
                logger.info("Processing PDF file")
                chunks = self._process_pdf(file_path)
                return [Document(page_content=chunk["content"], metadata=chunk["metadata"]) for chunk in chunks]
            
            else:
                logger.warning(f"Unsupported file type: {file_type}")
                raise ValueError(f"Unsupported file type: {file_type}")
            
        except Exception as e:
            logger.error(f"Error processing document: {str(e)}", exc_info=True)
            raise Exception(f"Error processing document: {str(e)}")

    def _process_pptx(self, file_path: str) -> List[Dict]:
        chunks = []
        loader = UnstructuredPowerPointLoader(file_path)
        documents = loader.load()
        
        for doc in documents:
            slide_number = doc.metadata.get("page_number", "unknown")
            
            slide_info = f"PowerPoint Slide {slide_number} Information:\n"
            slide_info += f"Content Type: {doc.metadata.get('content_type', 'unknown')}\n"
            slide_info += f"Content:\n{doc.page_content}\n"
            
            chunks.append({
                "content": slide_info,
                "metadata": {
                    "source": file_path,
                    "type": "ppt_slide_info",
                    "slide_number": str(slide_number),
                    "content_type": doc.metadata.get("content_type", "unknown")
                }
            })
            
            content_chunks = self.text_splitter.split_text(doc.page_content)
            for i, chunk in enumerate(content_chunks):
                chunks.append({
                    "content": chunk,
                    "metadata": {
                        "source": file_path,
                        "type": "ppt_content",
                        "slide_number": str(slide_number),
                        "chunk_index": str(i + 1),
                        "total_chunks": str(len(content_chunks))
                    }
                })
        
        prs = Presentation(file_path)
        presentation_info = self._get_presentation_info(prs, file_path)
        chunks.append(presentation_info)
        
        return chunks

    def _process_pdf(self, file_path: str) -> List[Dict]:
        chunks = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text = page.extract_text()
                
                # Split text into chunks
                text_chunks = self.text_splitter.split_text(text)
                
                for i, chunk in enumerate(text_chunks):
                    chunks.append({
                        "content": chunk,
                        "metadata": {
                            "source": file_path,
                            "type": "pdf",
                            "page": str(page_num + 1),
                            "chunk": str(i + 1)
                        }
                    })
        
        return chunks

    def _get_presentation_info(self, prs: Presentation, file_path: str) -> Dict:
        """Get presentation-level information."""
        presentation_info = f"PowerPoint Presentation Information:\n"
        presentation_info += f"Total Slides: {len(prs.slides)}\n"
        presentation_info += f"Slide Layouts Used: {', '.join(set(slide.slide_layout.name for slide in prs.slides))}\n"
        
        # Add slide titles and their sequence
        presentation_info += "\nSlide Sequence:\n"
        for i, slide in enumerate(prs.slides, 1):
            title = "No Title"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
                    break
            presentation_info += f"Slide {i}: {title}\n"
        
        return {
            "content": presentation_info,
            "metadata": {
                "source": file_path,
                "type": "ppt_presentation_info",
                "total_slides": str(len(prs.slides))
            }
        } 