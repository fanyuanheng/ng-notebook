import os
import tempfile
import pandas as pd
import pytest
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import PyPDF2
from fpdf import FPDF
from ng_notebook.services.document_processor import DocumentProcessor
from ng_notebook.services.sqlite_store import SQLiteStore

@pytest.fixture(scope="module")
def document_processor(tmp_path_factory):
    # Use a temp SQLiteStore for isolation
    sqlite_dir = tmp_path_factory.mktemp("sqlite_db")
    os.environ["SQLITE_DB_DIR"] = str(sqlite_dir)
    return DocumentProcessor(SQLiteStore())

def create_temp_csv(tmp_path):
    csv_path = tmp_path / "test.csv"
    df = pd.DataFrame({"col1": [1, 2], "col2": ["a", "b"]})
    df.to_csv(csv_path, index=False)
    return csv_path

def create_temp_excel(tmp_path):
    excel_path = tmp_path / "test.xlsx"
    df = pd.DataFrame({"x": [10, 20], "y": ["foo", "bar"]})
    df.to_excel(excel_path, index=False, sheet_name="Sheet1")
    return excel_path

def create_temp_pdf(tmp_path):
    pdf_path = tmp_path / "test.pdf"
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="Hello PDF!", ln=True)
    pdf.output(str(pdf_path))
    return pdf_path

def create_temp_pptx(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Hello PPTX!"
    prs.save(str(pptx_path))
    return pptx_path

def detect_mime_type(file_path):
    import magic
    return magic.from_file(str(file_path), mime=True)

def test_process_csv(document_processor, tmp_path):
    csv_path = create_temp_csv(tmp_path)
    file_type = detect_mime_type(csv_path)
    chunks = document_processor.process_document(str(csv_path), file_type)
    assert len(chunks) == 1
    assert "CSV File" in chunks[0].page_content
    assert chunks[0].metadata["file_type"] == file_type

def test_process_excel(document_processor, tmp_path):
    excel_path = create_temp_excel(tmp_path)
    file_type = detect_mime_type(excel_path)
    chunks = document_processor.process_document(str(excel_path), file_type)
    assert len(chunks) >= 1
    assert "Sheet:" in chunks[0].page_content
    assert chunks[0].metadata["file_type"] == file_type

def test_process_pdf(document_processor, tmp_path):
    pdf_path = create_temp_pdf(tmp_path)
    file_type = detect_mime_type(pdf_path)
    chunks = document_processor.process_document(str(pdf_path), file_type)
    assert len(chunks) >= 1
    assert "Hello PDF!" in " ".join([c.page_content for c in chunks])
    for chunk in chunks:
        assert chunk.metadata["type"] == "pdf"

def test_process_pptx(document_processor, tmp_path):
    pptx_path = create_temp_pptx(tmp_path)
    file_type = detect_mime_type(pptx_path)
    chunks = document_processor.process_document(str(pptx_path), file_type)
    assert len(chunks) >= 1
    found_slide = any("PowerPoint Slide" in c.page_content for c in chunks)
    found_presentation = any("PowerPoint Presentation Information" in c.page_content for c in chunks)
    assert found_slide or found_presentation
    for chunk in chunks:
        assert chunk.metadata["source"] == str(pptx_path) 